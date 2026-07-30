from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPTX_PATH = Path("docs/Healthcare-Lab-Architecture-OnePager.pptx")

NAVY = RGBColor(0x10, 0x2A, 0x43)
TEAL = RGBColor(0x00, 0xA8, 0x96)
BLUE = RGBColor(0x17, 0x6B, 0x87)
PALE = RGBColor(0xEA, 0xF4, 0xF7)
PALE_ALT = RGBColor(0xF7, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x24, 0x3B, 0x53)
MUTED = RGBColor(0x62, 0x7D, 0x8C)
BORDER = RGBColor(0xAF, 0xC9, 0xD3)
WARN = RGBColor(0xEF, 0x6A, 0x5B)


DISPLAY_ROWS = [
    ("紙速（Paper Speed）", "25 mm/s", "水平時間刻度；小格 1 mm = 0.04 秒，大格 5 mm = 0.2 秒"),
    ("電壓增益（Voltage Gain）", "10 mm/mV", "垂直振幅刻度；小格 1 mm = 0.1 mV，大格 5 mm = 0.5 mV"),
    ("顯示振幅範圍", "每導程 ±2 mV", "超過固定顯示範圍時顯示警告"),
    ("Baseline 校正", "開啟", "以各導程樣本中位數置中顯示"),
    ("導程順序", "I、II、III、aVR、aVL、aVF、V1–V6", "固定標準 12 導程順序"),
    ("顯示版面", "兩欄，每欄六列", "肢體導程與胸前導程分欄顯示"),
    ("SVG 尺寸", "1600 × 800 px", "瀏覽器顯示用的預設輸出尺寸"),
    ("兩欄間距", "0.25 秒", "兩欄波形之間的視覺間隔"),
    ("顯示定位", "展示用途", "名義刻度；不保證螢幕或列印後的實體毫米尺寸"),
]

LIMIT_ROWS = [
    ("支援的 SOP Class", "Twelve-lead ECG、General ECG Waveform Storage", "其他 DICOM 類型不顯示"),
    ("Waveform Sequence", "必須正好一組", "缺少或多組都會拒絕"),
    ("導程要求", "完整 12 個 SCPECG 導程", "不接受缺少、重複或無法識別的導程"),
    ("樣本格式", "Signed 16-bit（SS）", "其他 sample representation 暫不支援"),
    ("電壓單位", "UCUM V、mV、uV／µV", "解析後統一換算成 mV"),
    ("取樣頻率", "從 DICOM 讀取", "測試資料預期 1000 Hz，並非前端寫死"),
    ("每導程樣本上限", "10,000 samples", "避免過大的繪圖負載"),
    ("顯示時間上限", "10 秒", "目前以 10 秒 resting ECG 為目標"),
    ("WADO-RS Timeout", "預設 30 秒", "可由 dcm4chee profile 調整"),
    ("DICOM instance 上限", "32 MB", "超過大小限制即停止下載"),
    ("可設定紙速範圍", "1–100 mm/s", "目前 API 使用預設 25 mm/s"),
    ("可設定增益範圍", "1–100 mm/mV", "目前 API 使用預設 10 mm/mV"),
]


def set_text(shape, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, color, bold=bold, align=align)
    return shape


def add_rect(slide, x, y, w, h, fill, line=BORDER, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.7)
    return shape


def add_header(slide, kicker, title, subtitle):
    add_textbox(slide, 0.55, 0.25, 2.2, 0.28, kicker, 11, TEAL, bold=True)
    add_textbox(slide, 0.55, 0.56, 7.4, 0.48, title, 26, NAVY, bold=True)
    add_textbox(slide, 8.15, 0.5, 4.6, 0.42, subtitle, 11, MUTED, align=PP_ALIGN.RIGHT)
    bar = add_rect(slide, 0.55, 1.08, 12.2, 0.04, TEAL, TEAL)
    bar.line.fill.background()


def add_table_slide(prs, kicker, title, subtitle, rows, widths, compact=False):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for placeholder in tuple(slide.placeholders):
        placeholder.element.getparent().remove(placeholder.element)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    add_header(slide, kicker, title, subtitle)

    x, y, total_w = 0.55, 1.35, 12.2
    header_h = 0.46
    body_h = (5.67 - header_h) / len(rows)
    headers = ("參數／項目", "目前設定／限制", "說明")

    cursor = x
    for heading, width in zip(headers, widths):
        cell = add_rect(slide, cursor, y, width, header_h, NAVY, NAVY)
        set_text(cell, heading, 11, WHITE, bold=True)
        cursor += width

    font_size = 9.2 if compact else 10.2
    value_size = 9.2 if compact else 10.2
    for row_index, row in enumerate(rows):
        row_y = y + header_h + row_index * body_h
        cursor = x
        fill = WHITE if row_index % 2 == 0 else PALE_ALT
        for col_index, (value, width) in enumerate(zip(row, widths)):
            cell = add_rect(slide, cursor, row_y, width, body_h, fill)
            set_text(
                cell,
                value,
                value_size if col_index == 1 else font_size,
                BLUE if col_index == 1 else TEXT,
                bold=col_index in (0, 1),
            )
            cursor += width

    add_textbox(
        slide,
        0.58,
        7.12,
        11.55,
        0.22,
        "Healthcare Lab ECG Viewer｜Demonstration only — not for diagnostic use",
        8.5,
        WARN,
        bold=True,
    )
    add_textbox(slide, 12.2, 7.12, 0.55, 0.22, str(len(prs.slides)), 8.5, MUTED, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation(PPTX_PATH)
    add_table_slide(
        prs,
        "DICOM ECG VIEWER",
        "ECG 顯示參數",
        "固定校正基準與瀏覽器呈現設定",
        DISPLAY_ROWS,
        (2.7, 3.35, 6.15),
    )
    add_table_slide(
        prs,
        "DICOM ECG VIEWER",
        "DICOM 資料與系統限制",
        "解析契約、渲染邊界與 WADO-RS 保護",
        LIMIT_ROWS,
        (2.65, 4.0, 5.55),
        compact=True,
    )
    temp = PPTX_PATH.with_name(PPTX_PATH.stem + ".updated.pptx")
    prs.save(temp)
    temp.replace(PPTX_PATH)


if __name__ == "__main__":
    main()
